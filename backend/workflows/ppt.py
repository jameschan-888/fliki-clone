"""PPT to video workflow (fliki 缺失工作流 #2)

- POST /workflow-ppt {slides: [{title, content}], language?, title?} -> workflow_draft
- MVP: 接收前端解析后的 slides JSON, 每页一个 scene
- 后续 P1: 直接接 .pptx 上传
"""
import uuid
from pathlib import Path
from workflows import build_workflow_router


def _ppt_to_scenes(body, language):
    slides = body.get("slides") or []
    if not isinstance(slides, list) or not slides:
        return ([], "")
    scenes = []
    for slide in slides:
        if not isinstance(slide, dict):
            continue
        content = (slide.get("content") or slide.get("text") or "").strip()
        if not content:
            continue
        scenes.append({
            "title": (slide.get("title") or ("幻灯片 " + str(len(scenes) + 1)))[:200],
            "narration": content[:5000],
            "subtitle": content[:200],
        })
    src = ("PPT 共 " + str(len(scenes)) + " 页") if scenes else ""
    return (scenes, src)



def _parse_pptx(path):
    try:
        from pptx import Presentation
    except ImportError:
        return []
    try:
        pres = Presentation(path)
    except Exception:
        return []
    slides = []
    for slide in pres.slides:
        title = ""
        body = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                txt = "".join(run.text for run in para.runs).strip()
                if not txt:
                    continue
                if not title:
                    title = txt[:200]
                else:
                    body.append(txt)
        slides.append({"title": title or ("幻灯片 " + str(len(slides) + 1)), "content": " ".join(body)[:5000]})
    return slides

def _ppt_to_scenes_extended(body, language):
    if not body.get("slides"):
        path = body.get("pptx_path")
        if path:
            parsed = _parse_pptx(path)
            if parsed:
                body = dict(body)
                body["slides"] = parsed
    return _ppt_to_scenes(body, language)
def create_router(get_db):
    router = build_workflow_router(prefix="/workflow-ppt", tag="workflow-ppt", source_to_scenes=_ppt_to_scenes_extended, get_db=get_db, max_source_length=50000, source_label="slides")
    from fastapi import HTTPException, Request
    from auth_router import get_user_id_from_request
    @router.post("/upload")
    async def upload_pptx(request: Request):
        user_id = get_user_id_from_request(request)
        if not user_id:
            raise HTTPException(status_code=401, detail="Authentication required")
        form = await request.form()
        upload = form.get("pptx")
        if upload is None or not hasattr(upload, "read"):
            raise HTTPException(status_code=422, detail="pptx 文件不能为空")
        target_dir = Path("data/uploads/pptx")
        target_dir.mkdir(parents=True, exist_ok=True)
        target = target_dir / (str(user_id) + "_" + uuid.uuid4().hex + ".pptx")
        target.write_bytes(await upload.read())
        return {"pptx_path": str(target), "filename": target.name}
    return router